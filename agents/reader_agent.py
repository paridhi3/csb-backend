import os
import fitz  # PyMuPDF
from pptx import Presentation
from io import BytesIO

def process_case_study(file_name: str, file_bytes: bytes):
    """
    Extract text and metadata from an uploaded PDF or PPTX file in memory.
    Args:
        file_name: Name of the file (used to detect type)
        file_bytes: File content in bytes
    Returns:
        (case_text, extracted_metadata)
    """
    case_text = ""

    # PDF
    if file_name.lower().endswith(".pdf"):
        pdf_stream = BytesIO(file_bytes)
        with fitz.open(stream=pdf_stream, filetype="pdf") as doc:
            for page in doc:
                case_text += page.get_text()

    # PPTX
    elif file_name.lower().endswith(".pptx"):
        ppt_stream = BytesIO(file_bytes)
        prs = Presentation(ppt_stream)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    case_text += shape.text + "\n"

    else:
        raise ValueError("Unsupported file type. Only PDF and PPTX are allowed.")

    extracted_metadata = {
        "title": os.path.basename(file_name),
        "length": len(case_text),
    }

    return case_text, extracted_metadata
